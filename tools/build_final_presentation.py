from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/chevy_troubleshooter_final_presentation.pptx")
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


BG = RGBColor(9, 23, 39)
CARD = RGBColor(18, 37, 57)
CARD_ALT = RGBColor(244, 247, 250)
ACCENT = RGBColor(32, 168, 216)
ACCENT_2 = RGBColor(245, 196, 66)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(22, 33, 45)
MUTED = RGBColor(174, 191, 207)
SUCCESS = RGBColor(62, 180, 137)


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, page_no: int) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(SLIDE_HEIGHT_IN - 0.45),
        Inches(SLIDE_WIDTH_IN - 1.0),
        Inches(0.2),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Chevrolet Manual/FAQ GraphRAG Troubleshooter  |  {page_no:02d}"
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, section: str, title: str, subtitle: str | None = None) -> None:
    section_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(2.0), Inches(0.3))
    p = section_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = section
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(9.8), Inches(0.75))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.48), Inches(9.8), Inches(0.45))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Malgun Gothic"
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(2.02), Inches(1.15), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    *,
    fill_color: RGBColor = CARD,
    title_color: RGBColor = WHITE,
    body_color: RGBColor = WHITE,
    accent_color: RGBColor = ACCENT,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1.2)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    text = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.16), Inches(width - 0.44), Inches(height - 0.28))
    tf = text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = title_color

    for idx, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_before = Pt(2 if idx == 0 else 0)
        p.space_after = Pt(0)
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(12)
        p.font.color.rgb = body_color


def add_metric_card(slide, left: float, top: float, width: float, label: str, value: str, note: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.0)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = TEXT_DARK

    p = tf.add_paragraph()
    p.text = label
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(60, 77, 94)

    p = tf.add_paragraph()
    p.text = note
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(9.5)
    p.font.color.rgb = RGBColor(102, 120, 138)


def add_chip(slide, left: float, top: float, text: str, fill_color: RGBColor) -> None:
    width = max(1.0, 0.12 * len(text) + 0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = TEXT_DARK if fill_color != CARD else WHITE
    p.alignment = PP_ALIGN.CENTER


def add_flow(slide, labels: list[str], top: float) -> None:
    left = 0.62
    widths = [1.8, 1.65, 2.15, 1.65, 1.9, 1.85]
    for idx, label in enumerate(labels):
        width = widths[idx] if idx < len(widths) else 1.7
        shape_type = MSO_SHAPE.CHEVRON if idx < len(labels) - 1 else MSO_SHAPE.ROUNDED_RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(0.82))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT if idx % 2 == 0 else ACCENT_2
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = "Malgun Gothic"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        left += width - 0.08


def add_comparison_row(slide, top: float, issue: str, fix: str, result: str) -> None:
    columns = [
        (0.7, 2.6, issue, RGBColor(246, 234, 233), RGBColor(154, 47, 49)),
        (3.45, 4.1, fix, RGBColor(232, 241, 250), RGBColor(24, 92, 146)),
        (7.75, 4.85, result, RGBColor(232, 246, 238), RGBColor(34, 118, 74)),
    ]
    for left, width, text, fill_color, title_color in columns:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.82))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = "Malgun Gothic"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = title_color


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    prs.core_properties.title = "Chevrolet Manual/FAQ GraphRAG Troubleshooter"
    prs.core_properties.subject = "최종 발표용 프로젝트 설명"
    prs.core_properties.author = "OpenAI Codex"

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(8.7), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Chevrolet Manual/FAQ\nGraphRAG Troubleshooter"
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.65), Inches(2.75), Inches(7.8), Inches(1.0))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "페이지 중심 Retrieval로 답변 · 출처 · 대표 이미지 정합성을 개선한\n쉐보레 차량 매뉴얼/FAQ 기반 진단 시스템"
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(18)
    r.font.color.rgb = MUTED

    add_chip(slide, 0.65, 4.15, "Neo4j Manual Graph", ACCENT_2)
    add_chip(slide, 2.55, 4.15, "Chroma FAQ Store", ACCENT)
    add_chip(slide, 4.32, 4.15, "Cohere Rerank", ACCENT_2)
    add_chip(slide, 5.95, 4.15, "LangSmith Trace", ACCENT)
    add_chip(slide, 7.62, 4.15, "Page-Centered UX", ACCENT_2)

    quote = slide.shapes.add_textbox(Inches(8.75), Inches(1.1), Inches(3.9), Inches(4.4))
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "핵심 질문"
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2
    for text in [
        "왜 답변은 맞는데 이미지가 틀리는가?",
        "왜 FAQ가 있는데 PDF가 먼저 선택되는가?",
        "왜 Malibu와 말리부가 다른 결과를 만드는가?",
        "이 문제들을 어떻게 설명 가능하게 고쳤는가?",
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.bullet = True
    add_footer(slide, 1)

    # 2. Problems
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "01  Problem", "처음 부딪힌 핵심 문제", "단순한 텍스트 검색 성능보다, 근거 표현 방식과 라우팅 문제가 더 크게 드러났다.")
    add_card(slide, 0.65, 2.35, 3.0, 2.0, "대표 이미지 불일치", [
        "검색된 chunk와 표시 이미지가 다른 맥락처럼 보임",
        "사용자 입장에서는 '답변은 맞는데 출처가 틀려 보이는' 문제",
    ])
    add_card(slide, 3.85, 2.35, 3.0, 2.0, "FAQ가 PDF에 밀림", [
        "FAQ raw 데이터에 거의 동일 문장이 있어도 PDF manual이 먼저 선택됨",
        "하드 라우팅과 manual-first 병합이 원인",
    ])
    add_card(slide, 7.05, 2.35, 3.0, 2.0, "모델 alias 불안정", [
        "Malibu / 말리부 / ALL_NEW_말리부가 서로 다른 검색 결과를 만듦",
        "모델 인식과 DB 검색용 값이 분리되어 있지 않았음",
    ])
    add_card(slide, 10.25, 2.35, 2.45, 2.0, "출력 score 혼선", [
        "manual score와 FAQ score 의미가 달랐음",
        "UI의 '신뢰도 25%'는 설명 가능한 수치가 아니었음",
    ], fill_color=RGBColor(33, 53, 73))
    add_card(slide, 0.65, 4.7, 12.05, 1.45, "설계 기준 재정의", [
        "정답 문장 생성보다 먼저, source file · page · representative image가 같은 근거 단위를 가리키도록 구조를 바꿨다.",
        "이후 Manual과 FAQ를 둘 다 검색한 뒤 마지막에 선택하는 방향으로 라우팅을 재설계했다.",
    ], fill_color=RGBColor(16, 49, 73), accent_color=ACCENT_2)
    add_footer(slide, 2)

    # 3. Data and tech stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "02  Dataset & Stack", "데이터 규모와 기술 선택", "프로젝트 목표는 멀티소스 검색이 아니라 '설명 가능한 근거 선택'이었다.")
    add_metric_card(slide, 0.7, 2.2, 2.7, "PDF Manuals", "371", "쉐보레 매뉴얼 PDF 수")
    add_metric_card(slide, 3.6, 2.2, 2.7, "Pages", "11,532", "재적재 후 페이지 노드 수")
    add_metric_card(slide, 6.5, 2.2, 2.7, "Chunks", "22,343", "블록/문단 child chunk 수")
    add_metric_card(slide, 9.4, 2.2, 2.7, "FAQ", "144 / 10", "FAQ 항목 / 카테고리 수")
    add_card(slide, 0.7, 4.05, 3.0, 1.65, "Neo4j", [
        "Manual 구조: Brand → Model → Manual → Page → Chunk",
        "vector index + fulltext index + page metadata에 유리",
    ])
    add_card(slide, 3.95, 4.05, 2.7, 1.65, "Chroma", [
        "FAQ는 question/answer/category 중심",
        "간단한 vector store가 더 효율적",
    ])
    add_card(slide, 6.9, 4.05, 2.7, 1.65, "Models", [
        "LLM: gpt-4.1-mini",
        "Embedding: BAAI/bge-m3",
        "Reranker: Cohere rerank-v3.5",
    ])
    add_card(slide, 9.85, 4.05, 2.85, 1.65, "Observability", [
        "LangSmith trace / latency / cost 비교",
        "운영 단계에서 regression 분석 가능",
    ])
    add_footer(slide, 3)

    # 4. Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "03  End-to-End Flow", "최종 아키텍처 한 장 요약", "질의 해석과 source selection을 분리하여, Manual과 FAQ를 경쟁 후보로 다루는 구조로 정리했다.")
    add_flow(
        slide,
        ["User Query", "Guardrail", "Manual Retrieval", "FAQ Retrieval", "Source Selector", "Answer / UI"],
        2.45,
    )
    add_card(slide, 0.75, 3.7, 3.95, 2.0, "Manual Path", [
        "Neo4j vector search + fulltext search",
        "chunk fusion → page aggregation → Cohere rerank",
        "top page가 대표 이미지와 출처가 됨",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 4.95, 3.7, 3.1, 2.0, "FAQ Path", [
        "Chroma vector retrieval",
        "FAQ_INTENT_HINTS는 hard route가 아닌 soft hint",
        "FAQ도 항상 경쟁 후보로 남김",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 8.3, 3.7, 4.25, 2.0, "Selection & Output", [
        "FAQ top score vs manual top page score 비교",
        "FAQ가 이기면 FAQ만 전면 표시, top image는 비움",
        "Manual이 이기면 page-centered source + 대표 이미지 표시",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_footer(slide, 4)

    # 5. Ingestion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "04  Ingestion Design", "적재 파이프라인을 왜 이렇게 바꿨는가", "이미지보다 먼저 page retrieval 품질을 올리기 위해, 데이터 모델과 chunking부터 다시 설계했다.")
    add_flow(
        slide,
        ["Catalog", "PDF Parse", "Block Chunk", "Page Render", "Embedding", "Neo4j Upsert"],
        2.35,
    )
    add_card(slide, 0.75, 3.5, 3.85, 2.15, "고정 길이 chunk를 버린 이유", [
        "문맥 경계가 깨지고 page 단위 출처와 연결이 약했음",
        "사용자에게 보여줄 대표 이미지와 검색 단위가 어긋났음",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 4.8, 3.5, 3.75, 2.15, "블록/문단 chunk를 선택한 이유", [
        "child chunk에서 정밀도를 확보",
        "parent page를 최종 근거 단위로 안정화",
        "답변 · 출처 · 대표 이미지 정합성 개선",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_card(slide, 8.75, 3.5, 3.95, 2.15, "OCR / 3단 레이아웃 / 페이지 번호", [
        "텍스트 부족 시 OCR fallback",
        "3단 레이아웃은 읽기 순서 보정",
        "인쇄 페이지 번호 없으면 PDF index를 fallback으로 사용",
    ], fill_color=RGBColor(15, 43, 67))
    add_footer(slide, 5)

    # 6. Query understanding
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "05  Query Understanding", "모델명과 질의 의도를 먼저 정리하는 이유", "검색 실패의 상당수는 임베딩보다 모델 alias와 질의 라우팅 문제에서 발생했다.")
    add_card(slide, 0.75, 2.3, 5.2, 3.4, "Guardrail이 하는 일", [
        "타 브랜드 / 비자동차 요청 차단",
        "모델명 정규화",
        "model_candidates 확장",
        "fallback category 추론",
        "preferred_manual_types 추론",
        "FAQ_INTENT_HINTS 기반 FAQ hint 생성",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 6.15, 2.3, 6.0, 1.65, "모델 alias 해결 방식", [
        "입력: Malibu → 내부 family key: malibu → 검색 후보: [말리부, ALL_NEW_말리부, THE_NEW_말리부]",
        "DB 스키마를 다시 바꾸지 않고도, retrieval 직전 런타임 재매핑으로 세대명/영문/한글 차이를 흡수했다.",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_card(slide, 6.15, 4.25, 6.0, 1.45, "FAQ intent를 soft hint로 낮춘 이유", [
        "FAQ 키워드가 manual에도 존재할 수 있으므로 hard routing은 오탐에 취약했다.",
        "이제 FAQ hint는 '우선 검토' 신호이고, 최종 선택은 retrieval 결과를 보고 결정한다.",
    ], fill_color=RGBColor(15, 43, 67))
    add_footer(slide, 6)

    # 7. Manual retrieval
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "06  Manual Retrieval", "Manual 쪽은 왜 하이브리드 + 페이지 집계 구조인가", "정비/증상 질의는 semantic retrieval과 exact term retrieval이 모두 필요하고, 최종 출력 단위는 page여야 했다.")
    add_card(slide, 0.75, 2.25, 4.0, 3.65, "단계", [
        "1) bge-m3 dense vector search",
        "2) Neo4j fulltext search",
        "3) chunk-level score fusion",
        "4) 같은 page hit를 묶어 page score 계산",
        "5) Cohere rerank로 top page 재정렬",
        "6) query keyword overlap으로 최종 pruning",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 5.0, 2.25, 3.55, 3.65, "핵심 점수 아이디어", [
        "vector와 lexical을 둘 다 반영",
        "같은 page 내 다중 hit에는 support bonus",
        "manual_type이 query hint와 맞으면 bonus",
        "reranker는 최종 후보 집합에만 적용",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_card(slide, 8.8, 2.25, 3.9, 3.65, "왜 Page-Centered인가", [
        "사용자는 chunk id가 아니라 페이지와 이미지를 본다",
        "답변-출처-대표 이미지가 같은 객체를 가리키게 해야 신뢰가 생긴다",
        "top page를 대표 근거로 두면 UI와 설명이 훨씬 단순해진다",
    ], fill_color=RGBColor(15, 43, 67))
    add_footer(slide, 7)

    # 8. FAQ retrieval + selector
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "07  FAQ Retrieval & Source Selector", "FAQ를 '무조건 라우팅'하지 않고 경쟁 후보로 남긴 이유", "FAQ가 존재해도 manual이 먼저 보이던 문제를, 검색 실패가 아닌 병합 정책 문제로 보고 해결했다.")
    add_card(slide, 0.75, 2.25, 3.85, 3.45, "FAQ Retrieval", [
        "Chroma cosine similarity 기반 vector search",
        "FAQ JSON은 question / answer / category 구조로 단순 적재",
        "현재는 lexical search와 separate rerank는 아직 미적용",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 4.8, 2.25, 3.75, 3.45, "Selector 규칙", [
        "FAQ hint가 있어도 곧바로 FAQ로 확정하지 않음",
        "FAQ top score와 manual top page score 비교",
        "FAQ가 더 강하면 FAQ를 전면 표시하고 manual image는 숨김",
        "manual이 더 강하면 page-centered manual source를 사용",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_card(slide, 8.75, 2.25, 3.95, 3.45, "실제 해결 사례", [
        "질의: Bluetooth music ... Incoming/Outgoing Call ...",
        "FAQ raw 데이터에 동일 질문 존재",
        "기존: PDF infotainment page가 앞섬",
        "현재: FAQ source만 전면 표시, top image는 None",
    ], fill_color=RGBColor(15, 43, 67))
    add_footer(slide, 8)

    # 9. Output policy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "08  Output Policy", "대표 이미지와 score 출력을 왜 다시 설계했는가", "사용자가 이해할 수 없는 수치와 엇갈린 이미지는, 검색 품질이 좋아도 시스템 신뢰를 떨어뜨린다.")
    add_card(slide, 0.75, 2.3, 4.05, 3.35, "대표 이미지 정책", [
        "대표 이미지는 top manual page일 때만 표시",
        "FAQ가 최종 source이면 top image는 비움",
        "대표 이미지와 대표 출처는 반드시 같은 객체를 가리킴",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 5.0, 2.3, 3.55, 3.35, "score 분리", [
        "retrieval_score: 검색 내부 순위용",
        "rerank_score: Cohere 재정렬 점수",
        "relevance_label: 높음 / 보통 / 낮음",
        "confidence: 운영용 제어값",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_card(slide, 8.8, 2.3, 3.9, 3.35, "왜 이렇게 했는가", [
        "manual score와 FAQ score는 같은 척도가 아니었음",
        "UI의 '신뢰도 25%'는 설명하기 어려웠음",
        "이제 사용자에게는 관련도, 개발자에게는 retrieval/rerank를 분리해 제공",
    ], fill_color=RGBColor(15, 43, 67))
    add_footer(slide, 9)

    # 10. Troubleshooting log
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "09  Troubleshooting Log", "프로젝트 중 해결한 핵심 이슈", "각 이슈는 검색기 하나의 성능 문제가 아니라, 데이터 구조와 라우팅 정책의 문제였다.")
    add_card(slide, 0.68, 2.0, 2.55, 0.7, "문제", [], fill_color=RGBColor(95, 34, 38), title_color=WHITE, accent_color=RGBColor(187, 71, 74))
    add_card(slide, 3.43, 2.0, 4.0, 0.7, "해결 방법", [], fill_color=RGBColor(26, 66, 101), title_color=WHITE, accent_color=ACCENT)
    add_card(slide, 7.68, 2.0, 4.95, 0.7, "결과", [], fill_color=RGBColor(28, 90, 60), title_color=WHITE, accent_color=SUCCESS)
    add_comparison_row(slide, 2.9, "답변은 맞지만 대표 이미지가 엉뚱함", "Child 검색 → Parent Page 집계로 전환, 대표 이미지=top page PNG", "답변 · 출처 · 이미지가 같은 page 객체를 가리키도록 정렬")
    add_comparison_row(slide, 3.95, "Malibu와 말리부가 다른 결과를 만듦", "guardrail에서 family key 후 model_candidates 확장", "말리부 계열 manual만 안정적으로 좁혀서 검색")
    add_comparison_row(slide, 5.0, "FAQ가 진단 guardrail에 막혀 출력되지 않음", "브랜드 FAQ는 허용, FAQ path 추가", "포인트/계정/계약 등 FAQ 응답 정상화")
    add_comparison_row(slide, 6.05, "FAQ가 있어도 manual-first 병합에 밀림", "FAQ hint를 soft hint로 낮추고 source selector 추가", "Bluetooth/MyLink 같은 FAQ형 질문에서 FAQ가 전면 선택")
    add_footer(slide, 10)

    # 11. Observability & evaluation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "10  Observability & Evaluation", "왜 LangSmith와 다층 평가가 필요한가", "현재 시스템은 단순 answer F1보다 source/page/image alignment와 latency/cost가 더 중요하다.")
    add_card(slide, 0.75, 2.25, 3.95, 3.6, "LangSmith 도입 이유", [
        "workflow root trace + LangChain child run 연결",
        "질의별 latency, cost, regression 비교 가능",
        "retrieval / rerank / answer generation 병목 분석",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 4.95, 2.25, 3.75, 3.6, "주 평가 지표", [
        "routing accuracy",
        "source_file hit@k",
        "page hit@k",
        "FAQ/manual selection accuracy",
        "representative image alignment",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    add_card(slide, 8.95, 2.25, 3.7, 3.6, "보조 평가 지표", [
        "groundedness / hallucination rate",
        "RAGAS 또는 LLM-as-judge",
        "latency p50 / p95",
        "cost per query",
        "confidence calibration",
    ], fill_color=RGBColor(15, 43, 67))
    add_footer(slide, 11)

    # 12. Results and next steps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "11  Results & Next Step", "현재 상태와 이후 우선순위", "핵심 구조는 완성되었고, 이후 개선은 실패 케이스 누적과 FAQ 고도화가 중심이다.")
    add_card(slide, 0.75, 2.25, 4.0, 3.45, "현재 달성한 것", [
        "페이지 중심 retrieval 재구성",
        "대표 이미지-출처 page 정합성 확보",
        "모델 alias 검색 안정화",
        "FAQ soft hint + source selector 적용",
        "score 출력 분리",
        "LangSmith tracing 연결",
    ], fill_color=RGBColor(15, 43, 67), accent_color=SUCCESS)
    add_card(slide, 5.0, 2.25, 3.55, 3.45, "남아 있는 한계", [
        "FAQ는 아직 vector-only",
        "FAQ dedicated rerank 없음",
        "bbox가 없어 evidence crop 미적용",
        "confidence는 calibrated probability가 아님",
    ], fill_color=RGBColor(15, 43, 67))
    add_card(slide, 8.8, 2.25, 3.9, 3.45, "다음 우선순위", [
        "실패 케이스 누적",
        "필요 시 FAQ lexical search 추가",
        "그 다음 FAQ rerank 검토",
        "page/source 중심 평가 리포트 고도화",
    ], fill_color=RGBColor(15, 43, 67), accent_color=ACCENT_2)
    quote = slide.shapes.add_textbox(Inches(0.82), Inches(6.05), Inches(11.8), Inches(0.55))
    p = quote.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "결론: 이 프로젝트의 핵심은 더 복잡한 모델을 쓰는 것이 아니라, 올바른 근거 단위를 선택하고 그 근거를 사용자가 납득할 수 있게 보여주는 것이다."
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2
    add_footer(slide, 12)

    # 13. Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Q&A"
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(1.4), Inches(3.25), Inches(10.6), Inches(1.0))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Page-centered retrieval, FAQ soft hint, source selector, 그리고 LangSmith 운영 관점에서 질문 부탁드립니다."
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(18)
    r.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    add_chip(slide, 4.15, 4.6, "답변 · 출처 · 이미지 정합성", ACCENT_2)
    add_chip(slide, 6.65, 4.6, "Manual / FAQ 선택 기준", ACCENT)
    add_footer(slide, 13)

    return prs


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT_PATH)
    print(OUT_PATH.as_posix())


if __name__ == "__main__":
    main()
